"""Build the phase-feature research presentation (PPTX)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BLUE = RGBColor(0x00, 0x72, 0xB2)
VERM = RGBColor(0xD5, 0x5E, 0x00)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5B, 0x5B, 0x5B)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def setpar(p, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
           italic=False, space_after=8):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def accent_bar(s, color=VERM):
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def fig_slide(title, img, caption, cap_color=VERM):
    s = slide()
    accent_bar(s, cap_color)
    _, tf = box(s, Inches(0.55), Inches(0.28), Inches(12.4), Inches(0.9))
    setpar(tf.paragraphs[0], title, 26, INK, bold=True)
    # image centered
    from PIL import Image
    iw, ih = Image.open(img).size
    maxw, maxh = Inches(11.6), Inches(5.0)
    scale = min(maxw / iw, maxh / ih)
    w, h = int(iw * scale), int(ih * scale)
    left = int((SW - w) / 2)
    s.shapes.add_picture(img, left, Inches(1.25), width=w, height=h)
    _, cf = box(s, Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.7))
    setpar(cf.paragraphs[0], caption, 16, cap_color, bold=True,
           align=PP_ALIGN.CENTER)
    return s


# ---------------------------------------------------------------- 1 title
s = slide()
bg = s.background
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x10, 0x2A, 0x43)
_, tf = box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6))
setpar(tf.paragraphs[0],
       "A phase-evaluation feature for predicting hyperbolic volume from the "
       "Jones polynomial", 40, WHITE, bold=True)
p = tf.add_paragraph()
setpar(p, "Theory-guided feature engineering in an NSGA-II architecture search "
       "— and breaking the J₂ ceiling with the adjoint J₃",
       22, RGBColor(0xE8, 0x84, 0x3C), space_after=4)
_, tf2 = box(s, Inches(0.95), Inches(6.2), Inches(11.5), Inches(1.0))
setpar(tf2.paragraphs[0], "CatalinMoldova  ·  2026", 18,
       RGBColor(0xBFD, 0xBFD, 0xBFD) if False else RGBColor(0xC7, 0xD4, 0xE0))

# ---------------------------------------------------------------- 2 motivation
s = slide()
accent_bar(s)
_, tf = box(s, Inches(0.55), Inches(0.3), Inches(12.3), Inches(1.0))
setpar(tf.paragraphs[0], "The idea: let theory pick the feature", 30, INK, bold=True)
_, bf = box(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.4))
bullets = [
    ("Task", "Predict a knot's hyperbolic volume from its fundamental Jones "
     "polynomial J₂ (34 integer coefficients). 12,955 hyperbolic knots, "
     "≤13 crossings.", INK, False),
    ("Volume conjecture", "A specific evaluation of the colored Jones polynomial "
     "governs the volume. Jejjala et al. found |V(e^(3πi/4))| alone tracks "
     "the volume almost as well as a full network.", INK, False),
    ("The change", "Add one evaluation — V(q₀) at q₀ = e^(3πi/4) "
     "— as three features [Re, Im, |·|]. Feature vector 42 → 45. "
     "Nothing else changes: same search space, splits, and budget.", VERM, True),
    ("Test", "Re-run the identical NSGA-II architecture search with vs. without "
     "the feature. A clean A/B.", INK, False),
]
first = True
for tag, text, col, bold in bullets:
    p = bf.paragraphs[0] if first else bf.add_paragraph()
    first = False
    p.space_after = Pt(16)
    r = p.add_run(); r.text = f"{tag}.  "
    r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = col
    r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(20); r2.font.color.rgb = INK; r2.font.name = "Calibri"

# ---------------------------------------------------------------- 2b why learn it
s = slide()
accent_bar(s)
_, tf = box(s, Inches(0.55), Inches(0.3), Inches(12.3), Inches(1.0))
setpar(tf.paragraphs[0], "Why learn it: a cheap stand-in for an intractable limit",
       30, INK, bold=True)
why = [
    ("The catch", "The volume conjecture's actual claim only holds in the limit "
     "of infinitely high color N → ∞. Low color (J₂, J₂+J₃) is cheap to "
     "compute; that limit is computationally intractable to reach directly.",
     INK),
    ("The substitution", "Instead of chasing the limit, train on cheap "
     "low-color data paired with independently computed ground-truth volumes "
     "(SnapPy). The network learns, empirically, whatever relationship links "
     "the low-level polynomial coefficients to the high-level geometry.", VERM),
    ("What you get", "Learned pattern-recognition replaces the intractable "
     "exact computation: a single forward pass at inference extracts what "
     "would otherwise require an infeasible high-color limit — cheap at the "
     "low level, exploiting structure that only becomes meaningful at the "
     "higher level.", INK),
    ("One level higher", "Extension: predict the volume-correlated behavior of "
     "higher-color polynomials (J₄, J₅, …) directly from low-color data — "
     "avoiding their computation entirely, not merely the infinite-N limit.",
     BLUE),
]
_, bf = box(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.6))
first = True
for tag, text, col in why:
    p = bf.paragraphs[0] if first else bf.add_paragraph()
    first = False
    p.space_after = Pt(14)
    r = p.add_run(); r.text = f"{tag}.  "
    r.font.bold = True; r.font.size = Pt(19); r.font.color.rgb = col
    r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(19); r2.font.color.rgb = INK; r2.font.name = "Calibri"
s.notes_slide.notes_text_frame.text = (
    "Computing the Jones polynomial at low color (J2, or J2+J3) is cheap, but "
    "the volume conjecture's actual claim only holds in the limit of "
    "infinitely high color — a regime that's computationally intractable to "
    "reach directly. Rather than chasing that limit, this approach trains a "
    "network on cheap, low-color data paired with independently computed "
    "ground-truth volumes. The network learns, empirically, whatever complex "
    "relationship links these low-level polynomial coefficients to the "
    "high-level geometric quantity of volume. In effect, it substitutes "
    "learned pattern-recognition for the intractable exact computation, "
    "extracting at inference time in a single forward pass what would "
    "otherwise require an infeasible high-color limit. The result is a model "
    "that operates cheaply at the low level but effectively captures and "
    "exploits structure that only becomes meaningful at the higher level. A "
    "further extension would push this idea one level higher — training a "
    "network to predict the volume-correlated behavior of higher-color "
    "polynomials (J4, J5, …) directly from cheap low-color data, avoiding "
    "their direct computation entirely rather than merely avoiding the "
    "infinite-N limit.")

# ---------------------------------------------------------------- 3 correlation
fig_slide("A single phase evaluation nearly tracks the volume",
          "figures/feature_correlation.png",
          "|V(e^(3πi/4))| is the strongest single predictor in the set "
          "(r = 0.93 vs 0.80 for the previous best) — before any network.")

# ---------------------------------------------------------------- 3b phase sweep
fig_slide("Theory put the feature exactly where the signal peaks",
          "figures/phase_sweep.png",
          "Single-scalar correlation with volume is ≈ 0 at most evaluation "
          "angles and peaks right at q₀ = e^(3πi/4) (empirical max at 0.78π). "
          "The phase isn't tuned — the volume conjecture pointed to it.")

# ---------------------------------------------------------------- 4 A/B result
fig_slide("Result: same search budget, one added feature",
          "figures/ab_accuracy.png",
          "Ensemble test accuracy 97.14% → 97.31%; mean 96.93% → 97.11% "
          "(≈ 6% relative error reduction, ~3σ).")

# ---------------------------------------------------------------- 4b control
fig_slide("Control (measured): the gain is the phase, not just 3 more inputs",
          "figures/ab_ablation.png",
          "Same knee net [32,18,54,18], only the 3 extra columns vary: "
          "phase 97.31% > zeroed (J₂-only) 97.22% > Gaussian noise 97.11% "
          "ensemble. The improvement is the volume signal — noise columns "
          "even hurt slightly.")

# ---------------------------------------------------------------- 5 pareto
fig_slide("The whole Pareto front shifts down — but stays flat",
          "figures/pareto_overlay.png",
          "Lower error everywhere; ~180× more parameters barely helps. "
          "The bottleneck is the J₂ signal, not the model.")

# ---------------------------------------------------------------- 5b architecture
fig_slide("The best model: the phase-run's Pareto-knee architecture",
          "figures/architecture.png",
          "Found by the 45-feature (phase) search: 45→32→18→54→18→1 GELU MLP, "
          "4,101 params (lr 6.4e-3, dropout 0.006, L2 1e-6, batch 256) — "
          "test MRE 0.0289 ± 0.0004, ensemble 0.0269 (97.31%). Tiny, because "
          "the phase feature hands even a small net a near-direct line to the "
          "volume.", cap_color=BLUE)

# ---------------------------------------------------------------- 6 ceiling
s = slide()
accent_bar(s)
_, tf = box(s, Inches(0.55), Inches(0.3), Inches(12.3), Inches(1.0))
setpar(tf.paragraphs[0], "An information ceiling — and how to break it",
       30, INK, bold=True)

rows = [
    ("Input", "Dataset", "Accuracy (1 − MRE)"),
    ("J₂ (this work, 42 feat)", "≤13 cr., 12,955", "96.9%"),
    ("J₂ + phase (this work)", "≤13 cr., 12,955", "97.3%"),
    ("J₂ (Hughes et al. 2025)", "≤16 cr., 1.7M", "98.5% — best J₂"),
    ("J₃ adjoint (Hughes et al.)", "≤15 cr., 177k", "99.3% — highest"),
]
tbl_w, tbl_h = Inches(11.6), Inches(2.9)
gt = s.shapes.add_table(len(rows), 3, Inches(0.85), Inches(1.5),
                        tbl_w, tbl_h).table
gt.columns[0].width = Inches(4.6)
gt.columns[1].width = Inches(3.6)
gt.columns[2].width = Inches(3.4)
for ci, txt in enumerate(rows[0]):
    c = gt.cell(0, ci); c.text = txt
    c.fill.solid(); c.fill.fore_color.rgb = BLUE
    pr = c.text_frame.paragraphs[0]; pr.runs[0].font.bold = True
    pr.runs[0].font.color.rgb = WHITE; pr.runs[0].font.size = Pt(16)
for ri in range(1, len(rows)):
    highlight = ri == 2  # our phase result
    for ci, txt in enumerate(rows[ri]):
        c = gt.cell(ri, ci); c.text = txt
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0xFCE7, 0xFCE7, 0xFCE7) if False else (
            RGBColor(0xFB, 0xE7, 0xD8) if highlight else WHITE)
        r = c.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(15)
        r.font.bold = highlight or ci == 2
        r.font.color.rgb = VERM if highlight else INK

_, bf = box(s, Inches(0.7), Inches(4.7), Inches(12.0), Inches(2.4))
takeaways = [
    ("Takeaway", "A theory-motivated feature bought a real, budget-matched gain "
     "— for free, at every model size.", VERM),
    ("Ceiling", "But ~98.5% is set by the data: distinct knots share a J₂ yet "
     "differ ~3% in volume. No architecture resolves them.", INK),
    ("Next", "Measured in Part II: on a matched 1,419-knot set, the adjoint J₃ "
     "more than halves the error of J₂ — the ceiling is real and J₃ breaks it.",
     INK),
]
first = True
for tag, text, col in takeaways:
    p = bf.paragraphs[0] if first else bf.add_paragraph()
    first = False
    p.space_after = Pt(12)
    r = p.add_run(); r.text = f"{tag}.  "
    r.font.bold = True; r.font.size = Pt(19); r.font.color.rgb = col
    r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(19); r2.font.color.rgb = INK; r2.font.name = "Calibri"

# ---------------------------------------------------------------- 7 Part II methods
s = slide()
accent_bar(s, BLUE)
_, tf = box(s, Inches(0.55), Inches(0.3), Inches(12.3), Inches(1.0))
setpar(tf.paragraphs[0], "Part II (measured): the adjoint J₃ on a colored-Jones "
       "dataset", 30, INK, bold=True)

j3m = [
    ("Data", "New colored-Jones table: 1,426 knots ≤12 crossings with J₂ "
     "(fundamental) and J₃ (adjoint). Ground-truth volumes computed with SnapPy "
     "(ideal triangulation of the knot complement; exact by Mostow rigidity); "
     "7 non-hyperbolic torus knots dropped → 1,419 knots.", INK),
    ("Cleaning", "Mixed q/t variable conventions and framing-polluted exponent "
     "shifts → all features built on a shift-invariant canonical form. "
     "Validated: all 1,419 J₂ polynomials and volumes match the Part-I dataset.",
     INK),
    ("Phases re-derived", "The volume-correlated evaluation phase is swept per "
     "polynomial instead of assuming 3π/4 transfers: peaks at 13π/16 for J₂ "
     "(corr 0.943) and 5π/9 for J₃ (corr 0.959).", BLUE),
    ("Design", "Three NSGA-II searches at the identical budget (pop 40 / gens "
     "40 / 4+7 seeds / 500 epochs) on the same knots and split — J₂-only vs "
     "J₃-only vs J₂+J₃. Only the input polynomial differs.", VERM),
]
_, bf = box(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.6))
first = True
for tag, text, col in j3m:
    p = bf.paragraphs[0] if first else bf.add_paragraph()
    first = False
    p.space_after = Pt(14)
    r = p.add_run(); r.text = f"{tag}.  "
    r.font.bold = True; r.font.size = Pt(19); r.font.color.rgb = col
    r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(19); r2.font.color.rgb = INK; r2.font.name = "Calibri"

# ---------------------------------------------------------------- 8 J3 phase sweep
fig_slide("The volume-tracking phase moves with the color",
          "figures/phase_sweep_j2j3.png",
          "corr(log|V(e^(iθ))|, log vol) peaks at θ = 13π/16 for J₂ (0.943) "
          "and 5π/9 for J₃ (0.959) — the literature's 3π/4 is not the peak "
          "here. J₃'s single scalar is the strongest signal yet.",
          cap_color=BLUE)

# ---------------------------------------------------------------- 9 J2/J3 headline
fig_slide("J₃ breaks the J₂ ceiling — measured, matched budget",
          "figures/j2j3_accuracy.png",
          "Identical knots, split, and search budget: J₂ 94.84% → J₃ 97.67% → "
          "J₂+J₃ 97.92% ensemble accuracy. J₃ alone cuts the error of J₂ by "
          "2.2×; every knee stays under 2.7k params.")

# ---------------------------------------------------------------- 9b j2j3 arch
fig_slide("The best J₂+J₃ model: the Pareto-knee architecture",
          "figures/architecture_j2j3.png",
          "Found by the joined search: 68→27→27→1 GELU MLP, 2,647 params "
          "(lr 5.5e-3, dropout 0.007, L2 2.0e-6, batch 128) — test MRE "
          "0.0246 ± 0.0013, ensemble 0.0208 (97.92%). Weights: "
          "results_j2j3_j2j3/knee_model.pt.", cap_color=BLUE)

# ---------------------------------------------------------------- 10 J3 control
fig_slide("Control (measured): the J₂+J₃ gain rides on the phase columns",
          "figures/ab_ablation_j2j3.png",
          "Same [27, 27] GELU knee, only the 6 phase-evaluation columns vary: "
          "phase 97.92% > zeroed 97.56% > Gaussian noise 96.46% ensemble. "
          "The signal is real; noise columns actively hurt.")

# ---------------------------------------------------------------- 11 lit comparison
s = slide()
accent_bar(s, BLUE)
_, tf = box(s, Inches(0.55), Inches(0.3), Inches(12.3), Inches(1.0))
setpar(tf.paragraphs[0], "Context: accuracy per parameter against the literature",
       30, INK, bold=True)

rows = [
    ("Model", "Input", "Knots", "Params", "Test MRE"),
    ("Jejjala–Kar–Parrikar '19 — 2×100 MLP", "J₂", "313k ≤15 cr.", "~12,000",
     "2.45%"),
    ("Craven–Jejjala–Kar '21 — distilled 2×5 net", "|J₂(q₀)|", "313k ≤15 cr.",
     "~10²", "~4%"),
    ("Craven et al. '22 — 3×100 MLP", "J₂", "313k ≤15 cr.", "~22,000", "~3%"),
    ("Hughes et al. '25 — 5-layer MLP", "J₂", "177k ≤15 cr.", "~180,000",
     "1.85%"),
    ("Hughes et al. '25 — 5-layer MLP", "J₃", "177k ≤15 cr.", "~182,000",
     "0.62%"),
    ("Hughes et al. '25 — 5-layer MLP  (SOTA)", "J₂+J₃", "177k ≤15 cr.",
     "~187,000", "0.40%"),
    ("This work — NAS knee", "J₂", "1.4k ≤12 cr.", "879", "5.16%"),
    ("This work — NAS knee", "J₃", "1.4k ≤12 cr.", "1,095", "2.33%"),
    ("This work — NAS knee", "J₂+J₃", "1.4k ≤12 cr.", "2,647", "2.08%"),
]
gt = s.shapes.add_table(len(rows), 5, Inches(0.85), Inches(1.35),
                        Inches(11.6), Inches(4.1)).table
gt.columns[0].width = Inches(4.3)
gt.columns[1].width = Inches(1.5)
gt.columns[2].width = Inches(2.2)
gt.columns[3].width = Inches(1.8)
gt.columns[4].width = Inches(1.8)
for ci, txt in enumerate(rows[0]):
    c = gt.cell(0, ci); c.text = txt
    c.fill.solid(); c.fill.fore_color.rgb = BLUE
    pr = c.text_frame.paragraphs[0]; pr.runs[0].font.bold = True
    pr.runs[0].font.color.rgb = WHITE; pr.runs[0].font.size = Pt(15)
for ri in range(1, len(rows)):
    ours = ri >= 7
    highlight = ri == 9
    for ci, txt in enumerate(rows[ri]):
        c = gt.cell(ri, ci); c.text = txt
        c.fill.solid()
        c.fill.fore_color.rgb = (RGBColor(0xFB, 0xE7, 0xD8) if highlight
                                 else RGBColor(0xF2, 0xF2, 0xF2) if ours
                                 else WHITE)
        r = c.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(12)
        r.font.bold = highlight
        r.font.color.rgb = VERM if highlight else INK

_, bf = box(s, Inches(0.7), Inches(5.65), Inches(12.0), Inches(1.7))
notes = [
    ("Same trend, 70× fewer parameters", "The J₂ → J₃ → J₂+J₃ error ordering "
     "of Hughes et al. reproduces exactly, with ~70–200× smaller models found "
     "by the NAS.", VERM),
    ("Why absolute errors differ", "They train on ~125× more knots at higher "
     "crossing number — with J₂, accuracy is set by data volume, not "
     "architecture. Numbers across datasets are not directly comparable.",
     INK),
]
first = True
for tag, text, col in notes:
    p = bf.paragraphs[0] if first else bf.add_paragraph()
    first = False
    p.space_after = Pt(8)
    r = p.add_run(); r.text = f"{tag}.  "
    r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = col
    r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(15); r2.font.color.rgb = INK; r2.font.name = "Calibri"

# ---------------------------------------------------------------- 12 conclusions
s = slide()
accent_bar(s)
_, tf = box(s, Inches(0.55), Inches(0.3), Inches(12.3), Inches(1.0))
setpar(tf.paragraphs[0], "Conclusions and what comes next", 30, INK, bold=True)

conc = [
    ("Theory picks features", "The volume-conjecture phase evaluation is a "
     "near-free accuracy gain at every model size, on both datasets — and "
     "ablations confirm it is the signal, not the extra dimensions.", VERM),
    ("Information beats capacity", "Flat Pareto fronts and tiny (<3k-param) "
     "knees throughout: what limits accuracy is the invariant you feed in, "
     "not the network. J₂ → J₃ moves accuracy far more than any architecture "
     "change.", INK),
    ("Measured J₃ gain", "On a matched 1,419-knot set: J₂ 94.8% → J₃ 97.7% → "
     "J₂+J₃ 97.9%, reproducing the literature ordering with ~70× fewer "
     "parameters.", BLUE),
    ("Future work", "Scale J₃ to the ≤13-crossing census (12,955 knots) to "
     "close the data gap toward the published 99.3%; push the substitution "
     "one level higher — predict the volume-correlated behavior of J₄, J₅, … "
     "directly from cheap low-color data, avoiding their computation "
     "entirely; interpret the knee networks symbolically.", INK),
]
_, bf = box(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.6))
first = True
for tag, text, col in conc:
    p = bf.paragraphs[0] if first else bf.add_paragraph()
    first = False
    p.space_after = Pt(16)
    r = p.add_run(); r.text = f"{tag}.  "
    r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = col
    r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(20); r2.font.color.rgb = INK; r2.font.name = "Calibri"

prs.save("Jones_phase_feature.pptx")
print("saved Jones_phase_feature.pptx  slides:", len(prs.slides._sldIdLst))
